from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    text = ''
    # Load the PowerPoint presentation
    prs = Presentation(pptx_file)

    # Iterate through each slide in the presentation
    for slide in prs.slides:
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape is a table
            if shape.has_table:
                table = shape.table
                # Iterate through each row in the table
                for row in table.rows:
                    # Iterate through each cell in the row
                    for cell in row.cells:
                        # Extract text from the cell
                        text += cell.text + '\t'
                    text += '\n'  # Add a newline after each row

            # Check if the shape has text
            elif hasattr(shape, 'text'):
                text += shape.text + '\n'

            # Check if the shape is a group of shapes
            elif hasattr(shape, 'shapes'):
                # Iterate through each shape within the group
                for group_shape in shape.shapes:
                    # Check if the grouped shape has text
                    if hasattr(group_shape, 'text'):
                        text += group_shape.text + '\n'

    return text

# Path to your PowerPoint file
pptx_file = '/home/vtpldedpy/indic-backend/udaan-deploy-pipeline/pdf_upload/samgyaa.pptx'

# Extract text from the PowerPoint file
extracted_text = extract_text_from_pptx(pptx_file)

# Print the extracted text
print(extracted_text)
